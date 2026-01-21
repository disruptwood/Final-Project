from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Verified Circles: High-Trust Connections for Relocators"
    subtitle.text = "Tackling the Reliability Gap in Social Matching\nFinal Project"

    # Slide 2: The Problem Context
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Reliability Gap: Why matching isn't enough"
    content = slide.placeholders[1]
    content.text = (
        "• The Stats: 78% of relocators fail to build a core social circle within 12 months (Survey N=45).\n"
        "• The Pain: It's not about finding people. It's about the high failure rate of meetings.\n"
        "• Quote: \"I have 500 matches, but I stopped trying because everyone flakes.\" (Interview P3)\n"
        "• Core Issue: Low commitment mechanisms lead to a \"Market for Lemons\" (Low trust)."
    )

    # Slide 3: Methods + Evidence Map (Table)
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Methods + Evidence Map"
    
    shapes = slide.shapes
    rows, cols = 6, 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Method", "N / Source", "What it answers", "Key Limitation"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    
    # Check if we can style it simply by text manually (default is usually clear)
    
    data = [
        ["Market Analysis", "Cognitive Research (2024)", "\"Is this a real business?\"", "Macro data excludes niche apps."],
        ["Competitor Analysis", "Product Teardown", "\"Why do users leave?\"", "Functional analysis only."],
        ["Synthetic Interviews", "6 AI Personas (GPT-4)", "Hypothesis Generator", "Simulation only. Requires validation."],
        ["User Survey", "N=45 (Real Users - CSV)", "Validation of Isolation", "Biased to online distribution."],
        ["Scientific Theory", "Signaling Theory", "\"Why will verification work?\"", "Theoretical application."]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, item in enumerate(row_data):
            table.cell(row_idx+1, col_idx).text = item
            
    # Add Process Note below table
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = "Process Note: Synthetic interviews = hypothesis generator; validation = surveys + behavioral test."


    # Slide 4: Market Data
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The 'Social Discovery' Market is Growing vs Dating"
    
    # Content on left, Image on right
    # Actually, let's put bullets top, image bottom
    content = slide.placeholders[1]
    content.text = (
        "• Social Discovery Market: ~$2.5 Billion (2024) -> Proj. $7.6B by 2031 (CAGR 17%).\n"
        "• The Gap: Dating Apps earn $9B/year (Optimization).\n"
        "• The Opportunity: Users pay for Access to a high-reliability network."
    )
    
    # Add Image
    img_path = 'presentation/assets/survey_no_partner_behavior.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(3.5), height=Inches(3.5))

    # Slide 5: Competitor Analysis
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Competitor Analysis (The Trust & Commitment Gap)"
    
    content = slide.placeholders[1]
    content.text = (
        "• Bumble BFF: Optimizes for Visuals. Result: Dating fatigue.\n"
        "• Meetup: Optimizes for Topics. Result: High noise.\n"
        "• Nomad Apps: Optimizes for travelers. Result: High transience & Safety gaps.\n"
        "• Verified Circles: Optimizes for Trust & Commitment."
    )
    
    # Add Image
    img_path = 'presentation/assets/competitor_matrix.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), width=Inches(8))

    # Slide 6: Scientific Validation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Scientific Validation: Signaling Theory"
    content = slide.placeholders[1]
    content.text = (
        "1. Costly Signals: Trust is built when someone invests effort (Cost) that is hard to fake.\n"
        "   - Selfie = Cheap Signal (Low Trust)\n"
        "   - Career/Sport History = Costly Signal (High Trust)\n"
        "2. Homophily: We bond with those who share verified traits.\n"
        "3. Result: Verification acts as 'Social Collateral'."
    )

    # Slide 7: User Insights
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "User Insights (Validation)"
    content = slide.placeholders[1]
    content.text = (
        "• The Reliability Gap: 'People says maybe and vanish.' (Users crave structure).\n"
        "• The Trust Filter: Users currently stalk manually. We automate it.\n"
        "• Survey Validation (N=45): 44% 'Go Alone' when no partner found."
    )
    
    # Add Image
    img_path = 'presentation/assets/survey_no_partner_behavior.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(3.8), height=Inches(3.2))

    # Slide 8: Persona
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Persona: 'The Time-Constrained Professional'"
    content = slide.placeholders[1]
    content.text = (
        "• Name: Omer, 32 (Senior Dev, Relocated to NYC)\n"
        "• Pain: 'I have limited free time. I can't afford to waste an evening on a meetup where I don't click.'\n"
        "• Current Solution: Does nothing. (Risk of wasted time > Reward)\n"
        "• Needs: Efficiency, Reliability, Shared Context."
    )

    # Slide 9: Solution Strategy
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Solution: The '3-Layer Trust' System"
    content = slide.placeholders[1]
    content.text = (
        "Pivot: From 'Open Access' (Noise) to 'Verified Access' (Signal).\n\n"
        "1. Trust Layer: Identity Verification + Social Collateral. (Safety)\n"
        "2. Relevance Layer: Level-based Matching (e.g. Tennis 4.0). (Boredom)\n"
        "3. Commitment Layer: Reputation System + Event Deposit. (Flakiness)"
    )

    # Slide 10: Product Snapshot
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Product Snapshot & UVP"
    content = slide.placeholders[1]
    content.text = (
        "• UVP: 'The Reliability of a Club. The Convenience of an App.'\n"
        "• Tagline: 'No Flakes. No Fakes. Just Friends.'\n\n"
        "• Differentiation:\n"
        "   - Erasure of 'Swiping' -> Connect on Context.\n"
        "   - Erasure of 'Open Groups' -> Verified Peers.\n\n"
        "• Success Metrics:\n"
        "   - Primary: Show-up Rate (>85%)\n"
        "   - Secondary: Repeat Meet Rate"
    )

    # Slide 11: Quantitative Validation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Quantitative Validation Proposal (Next Steps)"
    content = slide.placeholders[1]
    content.text = (
        "• Behavioral Test (Smoke Test)\n"
        "• Landing Page: 'NYC Tech Runners - Verified Members Only.'\n"
        "• Metric: Conversion rate of visitors willing to link Strava.\n"
        "• Success Criteria: >15% conversion.\n"
        "  - Justification: Survey 2 showed 60% willingness so 15% is conservative."
    )

    # Slide 12: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = (
        "• The Problem: Loneliness is driven by Activity Isolation (44% go alone) and Flakiness.\n"
        "• The Solution: Verified Circles reduces no-show rates via Costly Signals.\n"
        "• Next Steps: Launch Smoke Test to validate conversion."
    )

    prs.save('presentation/deck.pptx')
    print("presentation.pptx created successfully.")

if __name__ == "__main__":
    create_presentation()
